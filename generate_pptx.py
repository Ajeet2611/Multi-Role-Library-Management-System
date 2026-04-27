"""
Library Management System Presentation Generator
12-slide professional .pptx in Hinglish
(Hindi for explanations + English for technical terms)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ===================== THEME =====================
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
NAVY_DARK = RGBColor(0x07, 0x14, 0x28)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
GOLD_LIGHT = RGBColor(0xF1, 0xC4, 0x0F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xE8, 0xEA, 0xED)
SOFT_GREY = RGBColor(0xB0, 0xB8, 0xC4)
CARD_BG = RGBColor(0x14, 0x2B, 0x4D)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_solid_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def add_bg(slide, color=NAVY):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_solid_fill(bg, color)
    no_line(bg)
    return bg


def add_gold_bar(slide, top=Inches(0), left=Inches(0), width=SLIDE_W, height=Inches(0.08)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_solid_fill(bar, GOLD)
    no_line(bar)
    return bar


def add_text(slide, left, top, width, height, text,
             font_size=18, bold=False, italic=False,
             color=WHITE, font="Calibri", align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items,
                font_size=16, color=WHITE, font="Calibri",
                bullet_color=GOLD, line_spacing=1.25):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet character
        b = p.add_run()
        b.text = "▸  "
        b.font.name = font
        b.font.size = Pt(font_size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        # text
        t = p.add_run()
        t.text = item
        t.font.name = font
        t.font.size = Pt(font_size)
        t.font.color.rgb = color
    return tb


def add_card(slide, left, top, width, height,
             fill=CARD_BG, border=GOLD, border_w=0.75):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.08
    set_solid_fill(card, fill)
    card.line.color.rgb = border
    card.line.width = Pt(border_w)
    return card


def add_title_block(slide, title_text, subtitle_text=None):
    """Standard title bar across the top of content slides."""
    add_gold_bar(slide, top=Inches(0))
    # Title
    add_text(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7),
             title_text, font_size=28, bold=True, color=WHITE,
             font="Calibri")
    if subtitle_text:
        add_text(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.45),
                 subtitle_text, font_size=15, italic=True, color=GOLD_LIGHT,
                 font="Calibri")
    # underline accent
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), Inches(1.35),
                                  Inches(1.2), Inches(0.05))
    set_solid_fill(line, GOLD)
    no_line(line)


def add_footer(slide, page_num, total=12):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
             "Multi-Role Library Management System  |  Ajeet Prasad",
             font_size=10, color=SOFT_GREY)
    add_text(slide, Inches(11.8), Inches(7.1), Inches(1.2), Inches(0.3),
             f"{page_num:02d} / {total:02d}",
             font_size=10, color=GOLD, align=PP_ALIGN.RIGHT)


def add_speaker_notes(slide, notes_text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = notes_text


def add_decoration(slide):
    """Decorative corner accents."""
    # Top-right gold triangle
    tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                 Inches(12.5), Inches(0),
                                 Inches(0.83), Inches(0.83))
    set_solid_fill(tri, GOLD)
    no_line(tri)
    tri.rotation = 90
    # Bottom-left small dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(0.2), Inches(7.1),
                                 Inches(0.2), Inches(0.2))
    set_solid_fill(dot, GOLD)
    no_line(dot)


# =====================================================================
# Build presentation
# =====================================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ============================ SLIDE 1: TITLE ============================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY_DARK)

# Decorative gold rings
ring1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(5), Inches(5))
ring1.fill.background()
ring1.line.color.rgb = GOLD
ring1.line.width = Pt(1.5)

ring2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(4.5), Inches(4), Inches(4))
ring2.fill.background()
ring2.line.color.rgb = GOLD
ring2.line.width = Pt(1.5)

# Top gold bar
add_gold_bar(s, top=Inches(0.3), left=Inches(5.2), width=Inches(2.9), height=Inches(0.06))

# Subtitle small
add_text(s, Inches(1), Inches(2.2), Inches(11.3), Inches(0.5),
         "PROJECT PRESENTATION  |  2026",
         font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER, font="Calibri")

# Main title
add_text(s, Inches(1), Inches(2.8), Inches(11.3), Inches(1.2),
         "Multi-Role Library Management System",
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

# Hindi tagline
add_text(s, Inches(1), Inches(4.0), Inches(11.3), Inches(0.6),
         "एक पूर्ण Production-Grade Desktop Application",
         font_size=22, italic=True, color=GOLD_LIGHT, align=PP_ALIGN.CENTER, font="Calibri")

# Tech tags
add_text(s, Inches(1), Inches(4.8), Inches(11.3), Inches(0.5),
         "Python  •  Tkinter  •  MySQL  •  bcrypt  •  Matplotlib",
         font_size=16, color=SOFT_GREY, align=PP_ALIGN.CENTER, font="Calibri")

# Author card
card = add_card(s, Inches(4.3), Inches(5.7), Inches(4.7), Inches(1.1),
                fill=CARD_BG, border=GOLD, border_w=1)
add_text(s, Inches(4.3), Inches(5.85), Inches(4.7), Inches(0.4),
         "Presented By", font_size=11, color=SOFT_GREY,
         align=PP_ALIGN.CENTER, font="Calibri")
add_text(s, Inches(4.3), Inches(6.15), Inches(4.7), Inches(0.5),
         "Ajeet Prasad", font_size=20, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font="Calibri")
add_text(s, Inches(4.3), Inches(6.55), Inches(4.7), Inches(0.4),
         "Senior Developer  |  Full-Stack Engineer", font_size=11,
         italic=True, color=GOLD_LIGHT, align=PP_ALIGN.CENTER, font="Calibri")

add_speaker_notes(s, """Pro-Tip: पहले 10 seconds सबसे important hain. Audience को बताइए कि ये सिर्फ एक college project नहीं — ये एक end-to-end engineered system hai jisme authentication, authorization, data layer, GUI और analytics — सब कुछ एक साथ integrate है।

Deep Dive: Mention कीजिए कि codebase में 15+ modules हैं और एक clean separation of concerns follow किया गया है (db / auth / security / GUI / business logic — सब अलग files में).""")


# ============================ SLIDE 2: PROBLEM ============================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_title_block(s, "Problem Statement",
                "आज भी 70% libraries manual registers पर depend करती हैं")
add_decoration(s)

problems = [
    ("📕", "Manual Tracking", "Books का manual record → lost data और human errors"),
    ("⏰", "Inconsistent Fines", "Late returns पर fine calculation में gaps"),
    ("🧑‍🤝‍🧑", "No Unified System", "Multiple roles के लिए कोई single platform नहीं"),
    ("🏢", "No Multi-Org Support", "College chains के लिए data isolation possible नहीं"),
    ("🔐", "Weak Security", "Plain-text passwords — serious security risk"),
    ("📩", "Poor Communication", "Student और Admin के बीच कोई proper channel नहीं"),
]

# 3x2 grid of problem cards
card_w = Inches(4.0)
card_h = Inches(1.55)
gap_x = Inches(0.25)
gap_y = Inches(0.25)
start_x = Inches(0.5)
start_y = Inches(1.7)

for i, (icon, title, desc) in enumerate(problems):
    row = i // 3
    col = i % 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    add_card(s, x, y, card_w, card_h, fill=CARD_BG, border=ACCENT_RED, border_w=0.75)
    # Icon
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(0.6), Inches(0.5),
             icon, font_size=24, color=WHITE, font="Segoe UI Emoji")
    # Title
    add_text(s, x + Inches(0.85), y + Inches(0.15), card_w - Inches(1.0), Inches(0.5),
             title, font_size=15, bold=True, color=GOLD, font="Calibri")
    # Description
    add_text(s, x + Inches(0.2), y + Inches(0.7), card_w - Inches(0.4), Inches(0.85),
             desc, font_size=12, color=LIGHT_GREY, font="Calibri")

# Bottom takeaway strip
strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.8))
strip.adjustments[0] = 0.3
set_solid_fill(strip, GOLD)
no_line(strip)
add_text(s, Inches(0.5), Inches(5.45), Inches(12.3), Inches(0.5),
         '"5000-student college में हर दिन 200 books — manually track करना almost impossible है"',
         font_size=16, bold=True, italic=True, color=NAVY_DARK,
         align=PP_ALIGN.CENTER, font="Calibri")

add_footer(s, 2)
add_speaker_notes(s, """Pro-Tip: Audience का emotional connection यहाँ बनता hai. Ek real example dijiye: Imagine kariye ek 5000-student college jahan har din 200 books issue hoti hain — manually track karna practically impossible है।

Deep Dive: Industry stat add कीजिए — NDLI report ke according, India में 60%+ academic libraries अभी bhi hybrid/manual systems पर हैं। यह credibility add करता है।

Animation suggestion: Staggered reveal — हर pain point एक-एक करके slide-in (left to right, 0.3s gap). Red accents pulse subtly to create urgency.""")


# ============================ SLIDE 3: SOLUTION ============================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_title_block(s, "Our Solution",
                "एक Unified, Secure, Multi-Role Library Operating System")
add_decoration(s)

solutions = [
    ("🔐", "Secure Login", "bcrypt-hashed passwords + role-based routing"),
    ("👥", "3-Tier Access", "Super Admin → Admin → User hierarchy"),
    ("📚", "Book Lifecycle", "Add → Issue → Return → Auto-fine — सब automated"),
    ("📊", "Live Analytics", "Real-time dashboard with Matplotlib charts"),
    ("💬", "Built-in Messaging", "User और Admin के बीच direct communication"),
    ("🏢", "Multi-Institution", "एक codebase, multiple organizations support"),
]

card_w = Inches(4.0)
card_h = Inches(1.55)
gap_x = Inches(0.25)
gap_y = Inches(0.25)
start_x = Inches(0.5)
start_y = Inches(1.7)

for i, (icon, title, desc) in enumerate(solutions):
    row = i // 3
    col = i % 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    add_card(s, x, y, card_w, card_h, fill=CARD_BG, border=GOLD, border_w=0.75)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(0.6), Inches(0.5),
             icon, font_size=24, color=WHITE, font="Segoe UI Emoji")
    add_text(s, x + Inches(0.85), y + Inches(0.15), card_w - Inches(1.0), Inches(0.5),
             title, font_size=15, bold=True, color=GOLD, font="Calibri")
    add_text(s, x + Inches(0.2), y + Inches(0.7), card_w - Inches(0.4), Inches(0.85),
             desc, font_size=12, color=LIGHT_GREY, font="Calibri")

# Bottom callout
strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.8))
strip.adjustments[0] = 0.3
set_solid_fill(strip, ACCENT_GREEN)
no_line(strip)
add_text(s, Inches(0.5), Inches(5.45), Inches(12.3), Inches(0.5),
         '"जो काम 6 अलग-अलग tools से होता था, अब एक ही Python application में"',
         font_size=16, bold=True, italic=True, color=WHITE,
         align=PP_ALIGN.CENTER, font="Calibri")

add_footer(s, 3)
add_speaker_notes(s, """Pro-Tip: यह slide problem से solution का aha moment है। एक hi line में बताइए: जो 6 अलग-अलग tools चाहिए थे, वो अब एक single Python application में हैं।

Deep Dive: Highlight कीजिए कि multi-institution support architecture-level decision है — हर query में InstitutionID filter है, जो data isolation guarantee करता है।

Animation: Mockup screens slide-in from depth. Feature icons orbit-rotate once around the mockup, then settle in hexagon. Each icon bounce-in with 0.2s stagger.""")


# ============================ SLIDE 4: ARCHITECTURE ============================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_title_block(s, "Solution Architecture",
                "Clean 3-Layer Design — Presentation, Logic, Data")
add_decoration(s)

# Three architecture layers
layers = [
    ("Layer 1 — Presentation (GUI)",
     "gui_login.py  •  gui_admin.py  •  gui_super_admin.py  •  gui_user.py",
     "Tkinter widgets के साथ Matplotlib embedded charts",
     ACCENT_BLUE),
    ("Layer 2 — Business Logic",
     "auth.py  •  permissions.py  •  messaging.py  •  admin_*.py modules",
     "Login flow, RBAC checks, और सारी business rules यहाँ implement",
     ACCENT_PURPLE),
    ("Layer 3 — Data Access",
     "db.py (Connection Factory)  •  MySQL Tables",
     "Users, Members, Books, Transactions, Messages, Institutions",
     ACCENT_GREEN),
]

layer_y = Inches(1.65)
layer_h = Inches(1.5)
layer_gap = Inches(0.2)

for i, (title, modules, desc, color) in enumerate(layers):
    y = layer_y + i * (layer_h + layer_gap)
    # Main card
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.8), y, Inches(11.7), layer_h)
    card.adjustments[0] = 0.08
    set_solid_fill(card, CARD_BG)
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    # Color tab on left
    tab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.8), y, Inches(0.25), layer_h)
    set_solid_fill(tab, color)
    no_line(tab)
    # Title
    add_text(s, Inches(1.2), y + Inches(0.15), Inches(11), Inches(0.45),
             title, font_size=18, bold=True, color=color, font="Calibri")
    # Modules (technical = English)
    add_text(s, Inches(1.2), y + Inches(0.6), Inches(11), Inches(0.4),
             modules, font_size=12, italic=True, color=GOLD_LIGHT, font="Consolas")
    # Description (Hindi)
    add_text(s, Inches(1.2), y + Inches(0.95), Inches(11), Inches(0.5),
             desc, font_size=13, color=LIGHT_GREY, font="Calibri")

# Flow arrows between layers (small triangles on the right)
for i in range(2):
    y = layer_y + (i + 1) * layer_h + i * layer_gap - Inches(0.05)
    arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                             Inches(11.8), y - Inches(0.05),
                             Inches(0.4), Inches(0.3))
    set_solid_fill(arr, GOLD)
    no_line(arr)

add_footer(s, 4)
add_speaker_notes(s, """Pro-Tip: Architecture slide में technical confidence दिखाना ज़रूरी है। बोलिए: हमने consciously separation of concerns follow किया है — GUI बदले तो logic को हाथ नहीं लगाना पड़ता, DB बदले तो UI safe रहती है।

Deep Dive: Mention कीजिए कि db.py एक factory pattern implement करता है — यह future में connection pooling या PostgreSQL migration को आसान बना देता है।

Animation: Build-up reveal — bottom layer first, then middle, then top. Data flow arrows pulse in a loop.""")


# ============================ SLIDE 5: TECH STACK ============================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_title_block(s, "Tech Stack",
                "Production-Ready Technologies — हर choice के पीछे एक reason")
add_decoration(s)

tech_items = [
    ("💻", "Python 3.10+", "Language", "Rapid development, huge ecosystem"),
    ("🖼️", "Tkinter", "GUI Framework", "Native, lightweight, zero dependency"),
    ("🗄️", "MySQL", "Database", "ACID compliance, reliable, free"),
    ("🔐", "bcrypt", "Security", "Industry-standard password hashing"),
    ("📊", "Matplotlib", "Analytics", "Embedded charts via FigureCanvasTkAgg"),
    ("📤", "openpyxl", "Excel Export", "Report generation for admins"),
    ("🖼️", "Pillow (PIL)", "Imaging", "Icon और asset rendering"),
    ("🧩", "MySQL Connector", "DB Driver", "Official Python ↔ MySQL bridge"),
]

# 4x2 grid
card_w = Inches(3.0)
card_h = Inches(1.4)
gap_x = Inches(0.15)
gap_y = Inches(0.2)
start_x = Inches(0.5)
start_y = Inches(1.7)

for i, (icon, name, category, reason) in enumerate(tech_items):
    row = i // 4
    col = i % 4
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    add_card(s, x, y, card_w, card_h, fill=CARD_BG, border=GOLD, border_w=0.75)
    # Icon top-left
    add_text(s, x + Inches(0.15), y + Inches(0.1), Inches(0.5), Inches(0.45),
             icon, font_size=22, color=WHITE, font="Segoe UI Emoji")
    # Category top-right
    add_text(s, x + Inches(0.7), y + Inches(0.15), card_w - Inches(0.85), Inches(0.3),
             category.upper(), font_size=9, bold=True, color=GOLD,
             align=PP_ALIGN.RIGHT, font="Calibri")
    # Name
    add_text(s, x + Inches(0.15), y + Inches(0.55), card_w - Inches(0.3), Inches(0.4),
             name, font_size=15, bold=True, color=WHITE, font="Calibri")
    # Reason (Hindi/Hinglish)
    add_text(s, x + Inches(0.15), y + Inches(0.9), card_w - Inches(0.3), Inches(0.5),
             reason, font_size=10, italic=True, color=LIGHT_GREY, font="Calibri")

# Bottom note
add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
         "💡 हर technology एक specific problem solve करती है — कोई भी choice random नहीं है",
         font_size=14, bold=True, italic=True, color=GOLD_LIGHT,
         align=PP_ALIGN.CENTER, font="Calibri")

add_footer(s, 5)
add_speaker_notes(s, """Pro-Tip: Tech stack slide पे सिर्फ logos मत दिखाओ — हर choice के पीछे reason बताओ। Example: MySQL चुना because relational data (Books → Transactions → Members) के लिए SQL JOINs सबसे efficient हैं।

Deep Dive: Mention कीजिए कि bcrypt deliberately slow है — यह brute-force attacks से protect करता है, जो plain SHA-256 नहीं कर सकता।

Animation: Logos fly-in from edges of the slide and snap into hexagon positions. Glow pulse on each logo as it lands.""")


# ============================ SLIDE 6: KEY FEATURES ============================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_title_block(s, "Key Features",
                "7 Killer Features जो इसे Production-Grade बनाते हैं")
add_decoration(s)

features = [
    ("🔐", "Secure Multi-Role Login",
     "bcrypt + dual-table lookup (Users + Members tables)"),
    ("📚", "Complete Book Management",
     "CRUD + ISBN tracking + multi-copy management"),
    ("🔄", "Smart Issue/Return Engine",
     "Auto-calculates fine ₹5/day on late returns"),
    ("📊", "Live Analytics Dashboard",
     "Pie chart (Active vs Inactive) + Bar chart (Roles)"),
    ("💬", "Two-Way Messaging",
     "User query → Admin reply → Status tracking (OPEN/CLOSED)"),
    ("📤", "Excel Export",
     "One-click report generation via openpyxl library"),
]

# Two columns
card_w = Inches(6.0)
card_h = Inches(1.5)
gap_y = Inches(0.15)
col_gap = Inches(0.3)

for i, (icon, title, desc) in enumerate(features):
    row = i // 2
    col = i % 2
    x = Inches(0.5) + col * (card_w + col_gap)
    y = Inches(1.65) + row * (card_h + gap_y)
    add_card(s, x, y, card_w, card_h, fill=CARD_BG, border=GOLD, border_w=0.75)
    # Big icon circle on left
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                x + Inches(0.2), y + Inches(0.3),
                                Inches(0.9), Inches(0.9))
    set_solid_fill(circle, NAVY_DARK)
    circle.line.color.rgb = GOLD
    circle.line.width = Pt(1)
    add_text(s, x + Inches(0.2), y + Inches(0.4), Inches(0.9), Inches(0.7),
             icon, font_size=26, color=WHITE,
             align=PP_ALIGN.CENTER, font="Segoe UI Emoji")
    # Title
    add_text(s, x + Inches(1.3), y + Inches(0.2), card_w - Inches(1.5), Inches(0.5),
             title, font_size=16, bold=True, color=GOLD, font="Calibri")
    # Desc
    add_text(s, x + Inches(1.3), y + Inches(0.7), card_w - Inches(1.5), Inches(0.7),
             desc, font_size=12, color=LIGHT_GREY, font="Calibri")

# Highlight strip
strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.45))
strip.adjustments[0] = 0.4
set_solid_fill(strip, GOLD)
no_line(strip)
add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.35),
         "✨  Bonus: Multi-Institution Isolation — हर query InstitutionID से filtered",
         font_size=13, bold=True, color=NAVY_DARK,
         align=PP_ALIGN.CENTER, font="Calibri")

add_footer(s, 6)
add_speaker_notes(s, """Pro-Tip: यहाँ सबसे impactful 1-2 features पे focus करो, बाकी का brief overview दो। मेरा suggestion — fine calculation और messaging system demo करो, क्योंकि ये real-world utility दिखाते हैं।

Deep Dive: Auto-fine logic interesting है — (date.today() - due_date).days * 5. साथ में बता सकते हो कि edge case (early return) भी handle होता है with 'if late > 0' check।

Animation: Cards flip-in one by one (like dealing playing cards). On hover, card tilts in 3D.""")


# ============================ SLIDE 7: DATABASE DESIGN ============================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_title_block(s, "Database Design",
                "Normalized Schema with 7 Core Tables (3NF)")
add_decoration(s)

tables = [
    ("Institutions", "Multi-org support", ACCENT_PURPLE, "PK: InstitutionID"),
    ("Users", "Admin/Super Admin accounts", ACCENT_BLUE, "FK → Institutions"),
    ("Members", "End-users (Students/Teachers)", ACCENT_BLUE, "FK → Institutions"),
    ("Books", "Inventory + copy tracking", ACCENT_GREEN, "FK → Institutions"),
    ("Transactions", "Issue/Return records", GOLD, "FK → Books, Members"),
    ("UserMessages", "Messaging system", ACCENT_RED, "Auto-created at runtime"),
    ("RolePermissions", "Granular RBAC matrix", ACCENT_PURPLE, "Many-to-Many"),
]

# Grid 4+3
card_w = Inches(3.0)
card_h = Inches(1.55)
gap = Inches(0.15)

# Row 1: 4 tables
for i in range(4):
    name, desc, color, fk = tables[i]
    x = Inches(0.5) + i * (card_w + gap)
    y = Inches(1.7)
    add_card(s, x, y, card_w, card_h, fill=CARD_BG, border=color, border_w=1)
    # Color header
    header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                x, y, card_w, Inches(0.35))
    set_solid_fill(header, color)
    no_line(header)
    add_text(s, x, y + Inches(0.03), card_w, Inches(0.3),
             name, font_size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font="Consolas")
    # Description
    add_text(s, x + Inches(0.15), y + Inches(0.5), card_w - Inches(0.3), Inches(0.5),
             desc, font_size=11, color=LIGHT_GREY, font="Calibri")
    # FK info
    add_text(s, x + Inches(0.15), y + Inches(1.05), card_w - Inches(0.3), Inches(0.4),
             fk, font_size=9, italic=True, color=GOLD_LIGHT, font="Consolas")

# Row 2: 3 tables (centered)
for i in range(3):
    name, desc, color, fk = tables[i + 4]
    x = Inches(2.1) + i * (card_w + gap)
    y = Inches(3.5)
    add_card(s, x, y, card_w, card_h, fill=CARD_BG, border=color, border_w=1)
    header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                x, y, card_w, Inches(0.35))
    set_solid_fill(header, color)
    no_line(header)
    add_text(s, x, y + Inches(0.03), card_w, Inches(0.3),
             name, font_size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font="Consolas")
    add_text(s, x + Inches(0.15), y + Inches(0.5), card_w - Inches(0.3), Inches(0.5),
             desc, font_size=11, color=LIGHT_GREY, font="Calibri")
    add_text(s, x + Inches(0.15), y + Inches(1.05), card_w - Inches(0.3), Inches(0.4),
             fk, font_size=9, italic=True, color=GOLD_LIGHT, font="Consolas")

# Key relationships strip
strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.3))
strip.adjustments[0] = 0.1
set_solid_fill(strip, NAVY_DARK)
strip.line.color.rgb = GOLD
strip.line.width = Pt(1)

add_text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(0.4),
         "🔗 Key Relationships",
         font_size=14, bold=True, color=GOLD, font="Calibri")
add_text(s, Inches(0.7), Inches(5.9), Inches(12), Inches(0.3),
         "▸  Books → Transactions → Members  (1 : N : 1)",
         font_size=12, color=WHITE, font="Calibri")
add_text(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.3),
         "▸  Institutions → Users / Members / Books  (1 : N for each)",
         font_size=12, color=WHITE, font="Calibri")
add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.3),
         "▸  हर query में InstitutionID filter — Multi-tenant data isolation guaranteed",
         font_size=12, italic=True, color=GOLD_LIGHT, font="Calibri")

add_footer(s, 7)
add_speaker_notes(s, """Pro-Tip: Audience को बता दो कि schema 3NF normalized है — कोई redundant data नहीं। Example: book का title सिर्फ Books table में, transactions सिर्फ BookID को reference करते हैं।

Deep Dive: UserMessages table runtime पे CREATE TABLE IF NOT EXISTS से बनती है — यह self-healing schema approach है, जो deployment को foolproof बनाता है।

Animation: Each table box fades in in dependency order (Institutions first, then Users/Members/Books, then Transactions). FK lines animate-draw between connected tables.""")


# ============================ SLIDE 8: CODE WORKFLOW ============================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_title_block(s, "Code Workflow",
                "End-to-End Request Lifecycle: Login → Dashboard")
add_decoration(s)

steps = [
    ("1", "GUI", "User opens gui_login.py", ACCENT_BLUE),
    ("2", "INPUT", "Credentials entered → do_login() triggered", ACCENT_BLUE),
    ("3", "QUERY", "auth.login() — Users table check", ACCENT_PURPLE),
    ("4", "FALLBACK", "Not found → Members table lookup", ACCENT_PURPLE),
    ("5", "VERIFY", "security.check_password() — bcrypt verify", ACCENT_RED),
    ("6", "ROUTE", "Returns user dict {Role, InstitutionID}", GOLD),
    ("7", "DASHBOARD", "Role-based routing → 3 possible dashboards", ACCENT_GREEN),
    ("8", "PERMS", "permissions.has_permission() loaded", ACCENT_GREEN),
]

# Vertical timeline on left, content on right
timeline_x = Inches(1.2)
content_x = Inches(2.0)
y_start = Inches(1.65)
step_h = Inches(0.62)

# Vertical line
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          timeline_x + Inches(0.13), y_start + Inches(0.1),
                          Inches(0.04), step_h * len(steps) - Inches(0.2))
set_solid_fill(line, GOLD)
no_line(line)

for i, (num, tag, desc, color) in enumerate(steps):
    y = y_start + i * step_h
    # Number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                timeline_x, y,
                                Inches(0.3), Inches(0.3))
    set_solid_fill(circle, color)
    circle.line.color.rgb = WHITE
    circle.line.width = Pt(1)
    add_text(s, timeline_x, y, Inches(0.3), Inches(0.3),
             num, font_size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font="Calibri")
    # Tag (English technical)
    add_text(s, content_x, y, Inches(1.5), Inches(0.3),
             tag, font_size=11, bold=True, color=color, font="Consolas")
    # Description
    add_text(s, content_x + Inches(1.6), y, Inches(9.5), Inches(0.4),
             desc, font_size=12, color=LIGHT_GREY, font="Calibri")

# Bottom: 3 dashboard outcomes
add_text(s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.3),
         "🎯 Final Routes:  SUPER_ADMIN → Super Dashboard   |   ADMIN → Admin Dashboard   |   USER → User Dashboard",
         font_size=11, italic=True, bold=True, color=GOLD_LIGHT,
         align=PP_ALIGN.CENTER, font="Calibri")

add_footer(s, 8)
add_speaker_notes(s, """Pro-Tip: इस slide पे एक live walk-through करो — actual code snippets दिखा सकते हो auth.py से। Audience को relate करने में मदद मिलेगी।

Deep Dive: Mention करो कि dual-table fallback (Users + Members) एक deliberate design choice है — Admins और end-users को अलग tables में रखने से permission management cleaner होता है, but login experience unified रहता है।

Animation: Flowchart draws itself sequentially. A glowing "request" particle travels along the flow showing real-time data movement.""")


# ============================ SLIDE 9: SECURITY ============================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_title_block(s, "Security Implementation",
                "Security First — कोई afterthought नहीं")
add_decoration(s)

sec_items = [
    ("🔐", "Password Security",
     "bcrypt + auto-generated salt",
     ["Adaptive cost factor — slow by design = brute-force resistant",
      "Plain passwords कभी store या log नहीं होते"]),
    ("👥", "Role-Based Access (RBAC)",
     "3-tier hierarchy + granular permissions",
     ["SUPER_ADMIN > ADMIN > USER hierarchy",
      "InstitutionRolePermissions table से fine-grained control"]),
    ("🛡️", "SQL Injection Prevention",
     "Parameterized queries everywhere",
     ["%s placeholders use होते हैं — कभी string concatenation नहीं",
      "MySQL Connector driver-level protection"]),
    ("🏢", "Multi-Tenancy Isolation",
     "InstitutionID filter on every query",
     ["Institution A का admin Institution B का data कभी नहीं देख सकता",
      "Architecture-level guarantee — application bug-proof"]),
]

card_w = Inches(6.0)
card_h = Inches(2.4)
gap = Inches(0.3)

for i, (icon, title, sub, points) in enumerate(sec_items):
    row = i // 2
    col = i % 2
    x = Inches(0.5) + col * (card_w + gap)
    y = Inches(1.65) + row * (card_h + Inches(0.15))
    add_card(s, x, y, card_w, card_h, fill=CARD_BG, border=GOLD, border_w=1)
    # Icon
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(0.6), Inches(0.55),
             icon, font_size=28, color=WHITE, font="Segoe UI Emoji")
    # Title
    add_text(s, x + Inches(0.9), y + Inches(0.15), card_w - Inches(1.1), Inches(0.4),
             title, font_size=16, bold=True, color=GOLD, font="Calibri")
    # Subtitle
    add_text(s, x + Inches(0.9), y + Inches(0.55), card_w - Inches(1.1), Inches(0.35),
             sub, font_size=11, italic=True, color=GOLD_LIGHT, font="Consolas")
    # Bullets
    for j, pt in enumerate(points):
        py = y + Inches(1.0) + j * Inches(0.55)
        add_text(s, x + Inches(0.3), py, Inches(0.2), Inches(0.4),
                 "▸", font_size=14, bold=True, color=GOLD, font="Calibri")
        add_text(s, x + Inches(0.55), py, card_w - Inches(0.7), Inches(0.5),
                 pt, font_size=11, color=LIGHT_GREY, font="Calibri")

add_footer(s, 9)
add_speaker_notes(s, """Pro-Tip: Security slide पे confidence दिखाना सबसे important है। बोलो: हमने OWASP Top 10 के principles को mind में रखके design किया — A02 (Cryptographic Failures), A03 (Injection), A07 (Authentication Failures) — सब address किए गए हैं।

Deep Dive: bcrypt vs SHA-256 का comparison करो — SHA-256 1 second में billions hashes कर सकता है, bcrypt deliberately ~100ms लेता है — यह attacker के लिए game over है।

Animation: Shield assembles layer-by-layer from outside-in. Each layer glows as it locks into place.""")


# ============================ SLIDE 10: CHALLENGES ============================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_title_block(s, "Challenges Overcome",
                "Real Engineering Problems जो हमने solve किए")
add_decoration(s)

challenges = [
    ("🔄", "Tkinter blocks main thread",
     "Lightweight queries + connection scoped per operation"),
    ("🏢", "Multi-tenant data leak risk",
     "Mandatory InstitutionID filter at query level"),
    ("🔐", "Plain → Hashed password migration",
     "One-time fix_passwords.py migration script"),
    ("📊", "Matplotlib in Tkinter integration",
     "Used FigureCanvasTkAgg backend (non-trivial)"),
    ("📩", "Schema bootstrap problem",
     "CREATE TABLE IF NOT EXISTS — self-healing schema"),
    ("🔁", "Dual login source (Users + Members)",
     "Fallback lookup pattern in auth.py"),
    ("💰", "Late fine edge cases",
     "if late > 0 guard with explicit messaging"),
    ("🎨", "Asset management for icons",
     "Centralized assets/ folder + Pillow rendering"),
]

# 2-column layout
card_w = Inches(6.0)
card_h = Inches(1.15)
gap = Inches(0.25)

for i, (icon, problem, solution) in enumerate(challenges):
    row = i // 2
    col = i % 2
    x = Inches(0.5) + col * (card_w + gap)
    y = Inches(1.65) + row * (card_h + Inches(0.12))
    add_card(s, x, y, card_w, card_h, fill=CARD_BG, border=ACCENT_GREEN, border_w=0.75)
    # Icon
    add_text(s, x + Inches(0.15), y + Inches(0.15), Inches(0.5), Inches(0.45),
             icon, font_size=20, color=WHITE, font="Segoe UI Emoji")
    # Problem (red-tinted)
    add_text(s, x + Inches(0.7), y + Inches(0.1), card_w - Inches(0.85), Inches(0.4),
             "⚠ " + problem, font_size=12, bold=True, color=ACCENT_RED, font="Calibri")
    # Solution (green-tinted)
    add_text(s, x + Inches(0.7), y + Inches(0.55), card_w - Inches(0.85), Inches(0.5),
             "✓ " + solution, font_size=11, color=ACCENT_GREEN, font="Calibri")

add_footer(s, 10)
add_speaker_notes(s, """Pro-Tip: Audience यहाँ engineer की problem-solving maturity judge करती है। सिर्फ 'क्या किया' मत बताओ — 'क्यों एक problem थी' और 'क्यों ये solution चुना' — दोनों explain करो।

Deep Dive: fix_passwords.py migration script का mention करना important — यह दिखाता है कि real-world systems में legacy data migration एक serious concern होता है, और तुमने handle किया।

Animation: Climber animates across the mountain range from left to right. Each peak lights up as climber passes.""")


# ============================ SLIDE 11: SCALABILITY ============================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_title_block(s, "Scalability & Future Roadmap",
                "Desktop App से Cloud-Scale Platform तक")
add_decoration(s)

# Current strengths box (left)
add_card(s, Inches(0.5), Inches(1.65), Inches(5.5), Inches(2.3),
         fill=CARD_BG, border=ACCENT_GREEN, border_w=1)
add_text(s, Inches(0.7), Inches(1.8), Inches(5.1), Inches(0.4),
         "✅ Current Scalability Strengths",
         font_size=15, bold=True, color=ACCENT_GREEN, font="Calibri")
strengths = [
    "Multi-institution architecture पहले से ready",
    "Stateless query layer — connection pooling आसान",
    "Modular code → MySQL को PostgreSQL/SQLite से swap easy",
    "Loose coupling — हर module independently deployable",
]
for i, item in enumerate(strengths):
    y = Inches(2.25) + i * Inches(0.4)
    add_text(s, Inches(0.8), y, Inches(0.2), Inches(0.3),
             "▸", font_size=13, bold=True, color=GOLD, font="Calibri")
    add_text(s, Inches(1.05), y, Inches(4.9), Inches(0.4),
             item, font_size=11, color=LIGHT_GREY, font="Calibri")

# Roadmap box (right)
add_card(s, Inches(6.3), Inches(1.65), Inches(6.5), Inches(5.3),
         fill=CARD_BG, border=GOLD, border_w=1)
add_text(s, Inches(6.5), Inches(1.8), Inches(6.1), Inches(0.4),
         "🚀 Roadmap (Next 6–12 Months)",
         font_size=15, bold=True, color=GOLD, font="Calibri")

roadmap = [
    ("🌐", "Web Version", "Flask/Django backend + React frontend"),
    ("📱", "Mobile App", "Flutter (iOS + Android)"),
    ("☁️", "Cloud Database", "PostgreSQL on AWS RDS / Replit DB"),
    ("🔔", "Notifications", "Email (SMTP) + SMS (Twilio) reminders"),
    ("🔍", "Advanced Search", "Full-text search via tsvector / Elasticsearch"),
    ("🤖", "AI Recommendations", "Collaborative filtering for book suggestions"),
    ("📊", "BI Dashboard", "Power BI / Metabase integration"),
]
for i, (icon, title, desc) in enumerate(roadmap):
    y = Inches(2.3) + i * Inches(0.62)
    add_text(s, Inches(6.55), y, Inches(0.4), Inches(0.4),
             icon, font_size=16, color=WHITE, font="Segoe UI Emoji")
    add_text(s, Inches(7.0), y, Inches(2.2), Inches(0.4),
             title, font_size=13, bold=True, color=GOLD_LIGHT, font="Calibri")
    add_text(s, Inches(9.2), y, Inches(3.5), Inches(0.4),
             desc, font_size=10, italic=True, color=LIGHT_GREY, font="Calibri")

# Bottom strengths continuation strip
strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.5), Inches(4.1), Inches(5.5), Inches(2.85))
strip.adjustments[0] = 0.05
set_solid_fill(strip, NAVY_DARK)
strip.line.color.rgb = GOLD
strip.line.width = Pt(1)
add_text(s, Inches(0.7), Inches(4.25), Inches(5.1), Inches(0.4),
         "💡 Vision",
         font_size=15, bold=True, color=GOLD, font="Calibri")
add_text(s, Inches(0.7), Inches(4.7), Inches(5.1), Inches(2.2),
         "यह सिर्फ एक desktop app नहीं — यह एक platform foundation है जिसे SaaS भी बनाया जा सकता है।\n\nCurrent schema पहले से multi-tenant ready है। SaaS conversion के लिए सिर्फ authentication layer को OAuth/JWT पे shift करना है, बाकी सब काम करेगा।",
         font_size=12, color=LIGHT_GREY, font="Calibri")

add_footer(s, 11)
add_speaker_notes(s, """Pro-Tip: Scalability slide पे vision दिखाना ज़रूरी है। बोलो: ये सिर्फ एक desktop app नहीं — ये एक platform foundation है जिसे हम SaaS भी बना सकते हैं।

Deep Dive: Mention करो कि current schema already multi-tenant ready है — SaaS conversion के लिए सिर्फ authentication layer को OAuth/JWT पे shift करना होगा, बाकी सब काम करेगा।

Animation: Timeline draws left-to-right like a progress bar. Each milestone flag plants itself with a small bounce.""")


# ============================ SLIDE 12: CONCLUSION ============================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY_DARK)

# Decorative rings
ring = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3), Inches(3), Inches(6), Inches(6))
ring.fill.background()
ring.line.color.rgb = GOLD
ring.line.width = Pt(1)
ring2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-2), Inches(6), Inches(6))
ring2.fill.background()
ring2.line.color.rgb = GOLD
ring2.line.width = Pt(1)

add_gold_bar(s, top=Inches(0.5), left=Inches(5.7), width=Inches(1.9), height=Inches(0.06))

# Big "Thank You"
add_text(s, Inches(0.5), Inches(0.8), Inches(12.3), Inches(1.2),
         "धन्यवाद",
         font_size=64, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER, font="Calibri")
add_text(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.5),
         "Thank You",
         font_size=22, italic=True, color=WHITE,
         align=PP_ALIGN.CENTER, font="Calibri")

# Why this stands out card
add_card(s, Inches(1.5), Inches(2.8), Inches(10.3), Inches(2.0),
         fill=CARD_BG, border=GOLD, border_w=1)
add_text(s, Inches(1.5), Inches(2.95), Inches(10.3), Inches(0.4),
         "✨ यह Project क्यों Stand Out करता है",
         font_size=16, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER, font="Calibri")
takeaways = [
    "✓  End-to-end ownership — UI, logic, DB, security — सब design किया",
    "✓  Production-grade decisions — bcrypt, RBAC, multi-tenancy",
    "✓  15+ modules, clean architecture, ~3000 lines of meaningful code",
    "✓  Real-world ready — कल किसी library में deploy हो सकता है",
]
for i, item in enumerate(takeaways):
    y = Inches(3.4) + i * Inches(0.32)
    add_text(s, Inches(2), y, Inches(9.3), Inches(0.35),
             item, font_size=12, color=WHITE, font="Calibri",
             align=PP_ALIGN.LEFT)

# Contact card
add_card(s, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.5),
         fill=CARD_BG, border=GOLD_LIGHT, border_w=1)
add_text(s, Inches(1.5), Inches(5.15), Inches(10.3), Inches(0.4),
         "📬 Connect With Me",
         font_size=14, bold=True, color=GOLD_LIGHT,
         align=PP_ALIGN.CENTER, font="Calibri")
add_text(s, Inches(1.5), Inches(5.55), Inches(10.3), Inches(0.4),
         "Ajeet Prasad  |  Senior Developer & Full-Stack Engineer",
         font_size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font="Calibri")
add_text(s, Inches(1.5), Inches(5.95), Inches(10.3), Inches(0.4),
         "📧  ajeetkumarbarh52@gmail.com    🔗  linkedin.com/in/ajeet-prasad-dev",
         font_size=12, color=LIGHT_GREY,
         align=PP_ALIGN.CENTER, font="Calibri")

# Final tagline
add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
         '"अगर कोई library मुझे आज ये system deploy करने को कहे — मैं कल से ready हूँ।"',
         font_size=14, italic=True, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER, font="Calibri")

add_speaker_notes(s, """Pro-Tip: Closing impactful बनाओ। Last 30 seconds में audience याद रखेगी — इसलिए one strong sentence से band करो। Suggested: अगर आज कोई library मुझे ये system deploy करने को कहे — मैं कल से ready हूँ।

Deep Dive: Q&A के लिए तैयार रहो। Common questions:
• Tkinter क्यों, web क्यों नहीं? (answer: rapid prototyping + offline-first)
• MySQL क्यों, MongoDB क्यों नहीं? (answer: relational data structure, ACID compliance)
• Production में deploy करने में कितना time लगेगा? (answer: schema + bcrypt हैं, सिर्फ deployment infra setup)

Animation: Thank You fades in with scale-up effect. QR codes slide-in from sides. Background particles drift slowly.""")


# ============================ SAVE ============================
output_path = "output/Library_Management_System_Presentation.pptx"
prs.save(output_path)
print(f"✓ Presentation generated: {output_path}")
print(f"✓ Total slides: {len(prs.slides)}")
